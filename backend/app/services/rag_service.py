from pathlib import Path
import hashlib
import logging
import shutil
import subprocess
import tempfile
from typing import List, Tuple

from ..config import Settings, resolved_documents_dir
from .openai_service import embed_documents, extract_text_from_image, generate_chat_completion
from .vector_store import VectorStore, load_vector_store

logger = logging.getLogger(__name__)

# 질문 키워드 → RAG 우선 폴더 (불량 > 규정 > 그 외 knowledge)
_FAILURE_KEYWORDS = (
    "불량",
    "불량품",
    "결함",
    "defect",
    "디펙",
    "양불",
    "수율",
    "불량률",
)
_RULES_KEYWORDS = (
    "규정",
    "내규",
    "사규",
    "취업규칙",
    "복무규정",
    "복무",
    "징계",
    "징계규정",
    "휴가규정",
    "인사규정",
    "윤리규정",
    "취업 규칙",
    "휴가 규정",
    "징계 절차",
)

# RAG 청크 크기(문자). 긴 문서는 여러 청크로 나눠 검색·임베딩한다.
RAG_CHUNK_CHARS = 2000

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".log",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf 미설치: PDF를 건너뜁니다.")
        return ""
    try:
        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                parts.append(t)
        return "\n\n".join(parts)
    except Exception:
        logger.exception("PDF 파싱 실패: %s", path.name)
        return ""


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl 미설치: Excel을 건너뜁니다.")
        return ""
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        lines: list[str] = []
        try:
            for sheet in wb.worksheets:
                lines.append(f"## {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(x.strip() for x in cells):
                        lines.append("\t".join(cells))
        finally:
            wb.close()
        return "\n".join(lines)
    except Exception:
        logger.exception("Excel 파싱 실패: %s", path.name)
        return ""


def _read_xls(path: Path) -> str:
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd 미설치: .xls를 건너뜁니다.")
        return ""
    try:
        book = xlrd.open_workbook(str(path))
        lines: list[str] = []
        for sheet in book.sheets():
            lines.append(f"## {sheet.name}")
            for r in range(sheet.nrows):
                row = sheet.row(r)
                cells = [str(c.value) if c.value != "" else "" for c in row]
                if any(x.strip() for x in cells):
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    except Exception:
        logger.exception("XLS 파싱 실패: %s", path.name)
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx 미설치: PPTX를 건너뜁니다.")
        return ""
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    for row in tbl.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            parts.append("\t".join(cells))
        return "\n\n".join(parts)
    except Exception:
        logger.exception("PPTX 파싱 실패: %s", path.name)
        return ""


def _read_ppt_libreoffice(path: Path) -> str:
    """
    레거시 .ppt는 python-pptx로 읽을 수 없음. LibreOffice가 있으면 txt로 변환해 텍스트를 얻는다.
    """
    bin_path = shutil.which("soffice") or shutil.which("libreoffice")
    if not bin_path:
        return ""
    try:
        with tempfile.TemporaryDirectory(prefix="samkwang_ppt_") as td:
            td_path = Path(td)
            src = path.resolve()
            subprocess.run(
                [
                    bin_path,
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    td,
                    str(src),
                ],
                check=False,
                timeout=120,
                capture_output=True,
            )
            out = td_path / (path.stem + ".txt")
            if not out.is_file():
                txts = sorted(td_path.glob("*.txt"))
                if len(txts) == 1:
                    out = txts[0]
                else:
                    return ""
            return out.read_text(encoding="utf-8", errors="replace")
    except Exception:
        logger.exception("LibreOffice로 .ppt 변환 실패: %s", path.name)
        return ""


def _read_document_content(path: Path) -> str:
    """
    문서 내용을 텍스트로 읽는다.
    지원되지 않는 바이너리 파일은 빈 문자열로 반환한다.
    """

    if not path.is_file():
        return ""

    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTENSIONS:
        return extract_text_from_image(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix in (".xlsx", ".xlsm"):
        return _read_xlsx(path)
    if suffix == ".xls":
        return _read_xls(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".ppt":
        text = _read_ppt_libreoffice(path)
        if not text.strip():
            logger.info(
                "레거시 .ppt 텍스트 추출 실패(빈 결과). LibreOffice 설치 또는 .pptx 변환을 권장: %s",
                path.name,
            )
        return text
    if suffix not in TEXT_EXTENSIONS:
        return ""

    raw = path.read_bytes()
    for encoding in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _build_doc_key(path: Path) -> str:
    return str(path.resolve())


def _rel_path_for_document(path: Path, documents_root: Path) -> str:
    try:
        rel = path.resolve().relative_to(documents_root.resolve())
    except ValueError:
        logger.warning("문서 경로가 documents_root 밖입니다: %s (root=%s)", path, documents_root)
        return ""
    return rel.as_posix().replace("\\", "/").lower()


def _build_hash(content: str) -> str:
    return hashlib.sha256(content.encode("utf-8")).hexdigest()


def _chunk_text(text: str, max_chars: int = RAG_CHUNK_CHARS) -> list[str]:
    t = text.strip()
    if not t:
        return []
    if len(t) <= max_chars:
        return [t]
    return [t[i : i + max_chars] for i in range(0, len(t), max_chars)]


def _clear_file_chunks_from_store(store: VectorStore, base_doc_key: str) -> None:
    """동일 파일의 이전 청크(또는 레거시 단일 doc_key)를 제거한다."""
    before = len(store._documents)
    store._documents = [
        d
        for d in store._documents
        if d.get("metadata", {}).get("base_doc_key") != base_doc_key
        and d.get("metadata", {}).get("doc_key") != base_doc_key
    ]
    if len(store._documents) != before:
        store.persist()


def _upsert_path_chunks(
    path: Path,
    content: str,
    settings: Settings,
    store: VectorStore,
    documents_root: Path,
) -> tuple[list[str], int]:
    """
    파일 내용을 청크로 저장하고 임베딩을 붙인다.
    반환: (chunk doc_key 목록, 청크 수)
    """
    base = _build_doc_key(path)
    stat = path.stat()
    rel_path = _rel_path_for_document(path, documents_root)
    chunks = _chunk_text(content)
    if not chunks:
        return [], 0

    keys: list[str] = []
    try:
        emb_list = embed_documents(chunks)
    except Exception:
        logger.exception("청크 임베딩 배치 실패: %s", path.name)
        emb_list = []

    for i, ch in enumerate(chunks):
        doc_key = f"{base}#chunk_{i}"
        keys.append(doc_key)
        emb = emb_list[i] if emb_list and i < len(emb_list) else None
        embedding = (
            emb
            if emb and any(abs(x) > 1e-12 for x in emb)
            else None
        )
        store.upsert_document(
            content=ch,
            metadata={
                "doc_key": doc_key,
                "base_doc_key": base,
                "chunk_index": i,
                "filename": path.name,
                "rel_path": rel_path,
                "size": stat.st_size,
                "modified_ns": stat.st_mtime_ns,
                "content_hash": _build_hash(ch),
            },
            embedding=embedding,
        )
    return keys, len(chunks)


def purge_stored_vectors_for_file(path: Path, settings: Settings) -> None:
    """디스크에서 파일을 지운 뒤에도 호출 가능. resolve() 기준 doc_key로 벡터 스토어 청크를 제거한다."""
    key = _build_doc_key(path)
    store = load_vector_store(settings.vector_db_path)
    _clear_file_chunks_from_store(store, key)


def ingest_document(path: Path, settings: Settings) -> None:
    content = _read_document_content(path)
    if not content.strip():
        logger.info("RAG 인덱싱 건너뜀 (미지원/빈 파일): %s", path.name)
        return

    documents_root = resolved_documents_dir(settings)
    store = load_vector_store(settings.vector_db_path)
    base = _build_doc_key(path)
    _clear_file_chunks_from_store(store, base)
    keys, n = _upsert_path_chunks(path, content, settings, store, documents_root)
    logger.info("RAG 인덱싱 %s: 청크 %s개 (keys=%s)", path.name, n, len(keys))


def sync_documents_folder(settings: Settings) -> dict:
    """
    document 폴더 전체를 스캔해 벡터 스토어를 동기화한다.
    """

    doc_dir = resolved_documents_dir(settings)
    doc_dir.mkdir(parents=True, exist_ok=True)
    store = load_vector_store(settings.vector_db_path)

    indexed_chunks = 0
    indexed_files = 0
    skipped = 0
    doc_keys: set[str] = set()

    for path in sorted(doc_dir.rglob("*")):
        if not path.is_file():
            continue
        content = _read_document_content(path)
        if not content.strip():
            skipped += 1
            continue
        keys, n_chunks = _upsert_path_chunks(path, content, settings, store, doc_dir)
        if not keys:
            skipped += 1
            continue
        doc_keys.update(keys)
        indexed_chunks += n_chunks
        indexed_files += 1

    store.prune_documents(valid_doc_keys=doc_keys)
    return {
        "indexed": indexed_chunks,
        "indexed_files": indexed_files,
        "skipped": skipped,
    }


def resolved_rag_folder_prefix(raw: str | None, settings: Settings) -> str | None:
    """
    클라이언트가 보낸 rag_folder 문자열을 document 루트 기준 안전한 상대 경로로 정규화한다.
    None·빈 문자열·비정상 경로면 None (문서 검색 생략).
    """
    if raw is None:
        return None
    s = str(raw).strip().replace("\\", "/").strip("/")
    if not s:
        return None
    if ".." in Path(s).parts:
        return None
    root = resolved_documents_dir(settings).resolve()
    try:
        target = (root / s).resolve()
        target.relative_to(root)
    except ValueError:
        return None
    return s


def rag_scope_rel_path_prefixes(query: str, settings: Settings) -> list[str]:
    """질문에 맞는 document 하위 폴더 접두사 1개를 반환(리스트)."""
    q = (query or "").strip()
    ql = q.lower()
    for kw in _FAILURE_KEYWORDS:
        if kw in q or kw.lower() in ql:
            return [settings.rag_failure_subdir]
    for kw in _RULES_KEYWORDS:
        if kw in q or kw.lower() in ql:
            return [settings.rag_rules_subdir]
    return [settings.rag_knowledge_subdir]


def retrieve_matches_for_chat(
    query: str,
    settings: Settings,
    k: int = 5,
    rag_folder: str | None = None,
) -> tuple[list[tuple[str, dict]], str]:
    """
    rag_folder 가 유효하면 해당 document 하위 폴더만 검색한다.
    미지정이면 벡터 검색을 하지 않고 빈 결과를 반환한다(일반 LLM 답변용).
    """
    store = load_vector_store(settings.vector_db_path)
    forced = resolved_rag_folder_prefix(rag_folder, settings)
    if forced is not None:
        scope_label = f"선택 문서 폴더: {forced}"
        prefixes = [forced]
        primary = store.similarity_search(query, k=k, rel_path_prefixes=prefixes)
        return primary[:k], scope_label

    # 폴더 미지정(도구「자동」): 벡터 검색·문서 컨텍스트 없이 LLM 일반 답변만 사용
    return [], "문서 검색 없음"


def answer_question(
    query: str,
    conversation_id: str | None,
    model: str,
    settings: Settings,
    rag_folder: str | None = None,
) -> Tuple[str, List[str]]:
    matches, scope_label = retrieve_matches_for_chat(query, settings, k=5, rag_folder=rag_folder)
    print(f"[RAG Service] 요청된 모델: {model}")
    if matches:
        context = "\n---\n".join(doc for doc, _meta in matches)
        prompt = (
            f"Conversation: {conversation_id or 'new'}\n\n"
            f"검색 우선 영역: {scope_label}\n\n"
            f"Context:\n{context}\n\nQuestion:\n{query}"
        )
    else:
        prompt = f"Conversation: {conversation_id or 'new'}\n\nQuestion:\n{query}"
    answer = generate_chat_completion(prompt, model=model)
    references = [meta.get("filename", "unknown") for _doc, meta in matches]
    return answer, references

